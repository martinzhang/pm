"""
AI Blueprint -- AI chat with SSE streaming, file content reading, image description
"""
import os, json, base64, time, requests
from datetime import date, datetime, timezone
from flask import Blueprint, request, jsonify, g, Response, stream_with_context
from models import get_db
from config import (
    MINIMAX_API_KEY, MINIMAX_BASE, MINIMAX_MODEL,
    UPLOAD_DIR, PHASES, PHASE_MAP,
    FILE_CONTENT_PER_FILE, FILE_CONTENT_TOTAL,
    TEXT_EXTS, IMAGE_EXTS, VLM_MAX_PER_REQUEST,
)

bp = Blueprint("ai", __name__)

PHASE_ORDER = [p[0] for p in PHASES]


# ── File content helpers ──

def _describe_image(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    mime_map = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
                ".gif": "image/gif", ".webp": "image/webp"}
    mime = mime_map.get(ext, "image/jpeg")
    try:
        with open(filepath, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("ascii")
        resp = requests.post(
            f"{MINIMAX_BASE}/chat/completions",
            headers={
                "Authorization": f"Bearer {MINIMAX_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": "MiniMax-VL-01",
                "messages": [{
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                        {"type": "text", "text": "Please describe this image in detail in Chinese. Include colors, objects, text content, layout, and any design elements you see. Be specific and thorough."},
                    ],
                }],
                "max_tokens": 1024,
            },
            timeout=30,
        )
        if resp.status_code == 200:
            data = resp.json()
            content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
            if "<think>" in content and "</think>" in content:
                content = content[content.index("</think>") + 8:].strip()
            return content if content else None
    except Exception:
        pass
    return None


def _read_file_content(filepath, original_name):
    ext = os.path.splitext(original_name)[1].lower()
    try:
        if ext in TEXT_EXTS:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                return f.read(FILE_CONTENT_PER_FILE)

        if ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(filepath)
            text = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text.append(t)
                if sum(len(s) for s in text) > FILE_CONTENT_PER_FILE:
                    break
            return "\n".join(text)[:FILE_CONTENT_PER_FILE] if text else None

        if ext == ".docx":
            from docx import Document
            doc = Document(filepath)
            text = "\n".join(p.text for p in doc.paragraphs if p.text)
            return text[:FILE_CONTENT_PER_FILE] if text else None

        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(filepath, read_only=True, data_only=True)
            text = []
            for ws in wb.worksheets:
                text.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) if c is not None else "" for c in row]
                    text.append("\t".join(vals))
                    if sum(len(s) for s in text) > FILE_CONTENT_PER_FILE:
                        break
            wb.close()
            return "\n".join(text)[:FILE_CONTENT_PER_FILE] if text else None

        if ext == ".xls":
            import xlrd
            wb = xlrd.open_workbook(filepath)
            text = []
            for ws in wb.sheets():
                text.append(f"[Sheet: {ws.name}]")
                for rx in range(ws.nrows):
                    vals = [str(ws.cell_value(rx, cx)) for cx in range(ws.ncols)]
                    text.append("\t".join(vals))
                    if sum(len(s) for s in text) > FILE_CONTENT_PER_FILE:
                        break
            return "\n".join(text)[:FILE_CONTENT_PER_FILE] if text else None

        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(filepath)
            text = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                if parts:
                    text.append(f"[Slide {i}] " + " | ".join(parts))
                if sum(len(s) for s in text) > FILE_CONTENT_PER_FILE:
                    break
            return "\n".join(text)[:FILE_CONTENT_PER_FILE] if text else None

    except Exception:
        return None
    return None


def _analyze_projects():
    conn = get_db()
    uid = g.user["id"]
    is_admin = g.user["role"] == "admin"
    if is_admin:
        projects = [dict(r) for r in conn.execute(
            "SELECT * FROM projects WHERE status='active' ORDER BY updated_at DESC"
        ).fetchall()]
    else:
        projects = [dict(r) for r in conn.execute(
            "SELECT * FROM projects WHERE status='active' AND (owner_id=? OR id IN (SELECT DISTINCT project_id FROM tasks WHERE assignee_id=?)) ORDER BY updated_at DESC",
            (uid, uid),
        ).fetchall()]
    visible_pids = {p["id"] for p in projects}
    all_tasks = [dict(r) for r in conn.execute(
        "SELECT t.*, p.name as project_name FROM tasks t "
        "JOIN projects p ON t.project_id=p.id WHERE p.status='active'"
    ).fetchall()]
    all_files = [dict(r) for r in conn.execute(
        "SELECT tf.*, t.name as task_name, p.name as project_name "
        "FROM task_files tf JOIN tasks t ON tf.task_id=t.id "
        "JOIN projects p ON t.project_id=p.id WHERE p.status='active' "
        "ORDER BY tf.created_at DESC"
    ).fetchall()]
    all_project_files = []
    try:
        all_project_files = [dict(r) for r in conn.execute(
            "SELECT pf.*, p.name as project_name "
            "FROM project_files pf JOIN projects p ON pf.project_id=p.id "
            "WHERE p.status='active' ORDER BY pf.created_at DESC"
        ).fetchall()]
    except Exception:
        all_project_files = []
    all_subtasks = [dict(r) for r in conn.execute(
        "SELECT s.*, t.name as task_name FROM subtasks s "
        "JOIN tasks t ON s.task_id=t.id "
        "JOIN projects p ON t.project_id=p.id WHERE p.status='active'"
    ).fetchall()]
    conn.close()
    if not is_admin:
        all_tasks = [t for t in all_tasks if t["project_id"] in visible_pids]
        visible_tids = {t["id"] for t in all_tasks}
        all_files = [f for f in all_files if f["task_id"] in visible_tids]
        all_subtasks = [s for s in all_subtasks if s["task_id"] in visible_tids]
        all_project_files = [f for f in all_project_files if f["project_id"] in visible_pids]

    files_by_task = {}
    for f in all_files:
        files_by_task.setdefault(f["task_id"], []).append(f)
    project_files_by_pid = {}
    for f in all_project_files:
        project_files_by_pid.setdefault(f["project_id"], []).append(f)
    subtasks_by_task = {}
    for s in all_subtasks:
        subtasks_by_task.setdefault(s["task_id"], []).append(s)

    today = date.today().isoformat()
    now_dt = datetime.now()

    overdue = []
    stalled = []
    unassigned = []
    blockers = []
    deadline_pressure = []
    phase_summary = {p: 0 for p in PHASE_ORDER}

    task_by_id = {t["id"]: t for t in all_tasks}

    for t in all_tasks:
        if t["end_date"] and t["end_date"] < today and t["progress"] < 100:
            overdue.append(t)
        if t["updated_at"] and t["progress"] < 100:
            try:
                last = datetime.fromisoformat(t["updated_at"])
                if (now_dt - last).days >= 3:
                    stalled.append(t)
            except Exception:
                pass
        if not t["assignee_name"] and t["progress"] < 100:
            unassigned.append(t)
        if t["phase"] in phase_summary:
            phase_summary[t["phase"]] += 1
        if t["depends_on"] and t["progress"] < 100:
            try:
                deps = [int(x.strip()) for x in str(t["depends_on"]).split(",") if x.strip()]
                for dep_id in deps:
                    dep = task_by_id.get(dep_id)
                    if dep and dep["progress"] < 100:
                        blockers.append({
                            "blocked_task": t["name"],
                            "blocked_by": dep["name"],
                            "blocker_progress": dep["progress"],
                            "project": t["project_name"],
                        })
            except Exception:
                pass

    for p in projects:
        if p["deadline"]:
            try:
                dl = date.fromisoformat(p["deadline"])
                days_left = (dl - date.today()).days
                ptasks = [t for t in all_tasks if t["project_id"] == p["id"]]
                avg = sum(t["progress"] for t in ptasks) / len(ptasks) if ptasks else 0
                if days_left <= 7 and avg < 80:
                    deadline_pressure.append({
                        "project": p["name"],
                        "days_left": days_left,
                        "avg_progress": round(avg),
                        "task_count": len(ptasks),
                    })
            except Exception:
                pass

    lines = []
    lines.append("=== 项目总览 ===")
    lines.append(f"活跃项目数: {len(projects)}")
    lines.append(f"总任务数: {len(all_tasks)}")
    lines.append(f"已逾期: {len(overdue)}, 停滞: {len(stalled)}, 未分配: {len(unassigned)}")
    lines.append("")

    for p in projects:
        ptasks = [t for t in all_tasks if t["project_id"] == p["id"]]
        done = sum(1 for t in ptasks if t["progress"] == 100)
        lines.append(f"项目: {p['name']} | 状态: {p['status']} | 截止: {p.get('deadline','未设定')}")
        lines.append(f"  任务 {done}/{len(ptasks)} 完成")
        pfs = project_files_by_pid.get(p["id"], [])
        if pfs:
            lines.append(f"  项目附件 ({len(pfs)}): " + ", ".join(
                f"{f['original_name']} ({f.get('uploaded_by_name') or ''}, {(f.get('created_at') or '')[:10]})" for f in pfs
            ))
        for t in ptasks:
            dep_str = f" [依赖: {t['depends_on']}]" if t["depends_on"] else ""
            lines.append(
                f"  - {t['name']} | 阶段:{PHASE_MAP.get(t['phase'],t['phase'])} "
                f"| 进度:{t['progress']}% | 负责:{t['assignee_name'] or '未分配'} "
                f"| 截止:{t['end_date'] or '无'}{dep_str}"
            )
            subs = subtasks_by_task.get(t["id"], [])
            if subs:
                done_s = sum(1 for s in subs if s["is_done"])
                lines.append(f"    子任务 ({done_s}/{len(subs)}完成): " + ", ".join(
                    ("[v]" if s["is_done"] else "[ ]") + s["content"] for s in subs
                ))
            tfiles = files_by_task.get(t["id"], [])
            if tfiles:
                lines.append(f"    附件 ({len(tfiles)}): " + ", ".join(
                    f"{f['original_name']} ({f['uploaded_by_name']}, {f['created_at'][:10]})" for f in tfiles
                ))
        lines.append("")

    if overdue:
        lines.append("=== 逾期任务 ===")
        for t in overdue:
            lines.append(f"  - {t['name']} ({t['project_name']}) 截止:{t['end_date']} 进度:{t['progress']}%")
        lines.append("")

    if stalled:
        lines.append("=== 停滞任务 (3天+无更新) ===")
        for t in stalled:
            lines.append(f"  - {t['name']} ({t['project_name']}) 上次更新:{t['updated_at'][:10]}")
        lines.append("")

    if blockers:
        lines.append("=== 依赖阻塞 ===")
        for b in blockers:
            lines.append(f"  - \"{b['blocked_task']}\" 被 \"{b['blocked_by']}\" 阻塞 (阻塞方进度:{b['blocker_progress']}%)")
        lines.append("")

    if deadline_pressure:
        lines.append("=== 截止日压力 ===")
        for d in deadline_pressure:
            lines.append(f"  - {d['project']}: 还剩{d['days_left']}天, 平均进度{d['avg_progress']}%")
        lines.append("")

    if unassigned:
        lines.append("=== 未分配任务 ===")
        for t in unassigned:
            lines.append(f"  - {t['name']} ({t['project_name']})")

    lines.append("")
    lines.append(f"阶段分布: {', '.join(PHASE_MAP.get(k,k)+':'+str(v) for k,v in phase_summary.items())}")

    total_content_len = 0
    file_contents = []
    vlm_calls = 0
    for f in all_files:
        if total_content_len >= FILE_CONTENT_TOTAL:
            break
        filepath = os.path.join(UPLOAD_DIR, str(f["task_id"]), f["filename"])
        if not os.path.exists(filepath):
            continue
        ext = os.path.splitext(f["original_name"])[1].lower()

        if ext in IMAGE_EXTS:
            desc = f.get("description") or ""
            if not desc and vlm_calls < VLM_MAX_PER_REQUEST:
                desc = _describe_image(filepath) or ""
                if desc:
                    vlm_calls += 1
                    try:
                        conn2 = get_db()
                        conn2.execute("UPDATE task_files SET description=? WHERE id=?", (desc, f["id"]))
                        conn2.commit()
                        conn2.close()
                    except Exception:
                        pass
            if desc:
                total_content_len += len(desc)
                file_contents.append(
                    f"--- 图片: {f['original_name']} (任务: {f['task_name']}, 项目: {f['project_name']}) ---\n[视觉描述] {desc}"
                )
            continue

        content = _read_file_content(filepath, f["original_name"])
        if content:
            remaining = FILE_CONTENT_TOTAL - total_content_len
            content = content[:remaining]
            total_content_len += len(content)
            truncated = " (内容已截断)" if len(content) >= FILE_CONTENT_PER_FILE or len(content) >= remaining else ""
            file_contents.append(
                f"--- 文件: {f['original_name']} (任务: {f['task_name']}, 项目: {f['project_name']}){truncated} ---\n{content}"
            )

    if file_contents:
        lines.append("")
        lines.append("=== 附件文件内容 ===")
        for fc in file_contents:
            lines.append(fc)
            lines.append("")

    # Project-level files
    proj_file_contents = []
    for f in all_project_files:
        if total_content_len >= FILE_CONTENT_TOTAL:
            break
        filepath = os.path.join(UPLOAD_DIR, f"project_{f['project_id']}", f["filename"])
        if not os.path.exists(filepath):
            continue
        ext = os.path.splitext(f["original_name"])[1].lower()

        if ext in IMAGE_EXTS:
            desc = f.get("description") or ""
            if not desc and vlm_calls < VLM_MAX_PER_REQUEST:
                desc = _describe_image(filepath) or ""
                if desc:
                    vlm_calls += 1
                    try:
                        conn2 = get_db()
                        conn2.execute("UPDATE project_files SET description=? WHERE id=?", (desc, f["id"]))
                        conn2.commit()
                        conn2.close()
                    except Exception:
                        pass
            if desc:
                total_content_len += len(desc)
                proj_file_contents.append(
                    f"--- 项目图片: {f['original_name']} (项目: {f['project_name']}) ---\n[视觉描述] {desc}"
                )
            continue

        content = _read_file_content(filepath, f["original_name"])
        if content:
            remaining = FILE_CONTENT_TOTAL - total_content_len
            content = content[:remaining]
            total_content_len += len(content)
            truncated = " (内容已截断)" if len(content) >= FILE_CONTENT_PER_FILE or len(content) >= remaining else ""
            proj_file_contents.append(
                f"--- 项目文件: {f['original_name']} (项目: {f['project_name']}){truncated} ---\n{content}"
            )

    if proj_file_contents:
        lines.append("")
        lines.append("=== 项目附件文件内容 ===")
        for fc in proj_file_contents:
            lines.append(fc)
            lines.append("")

    return "\n".join(lines)


SYSTEM_PROMPT = """你是奈娃咖啡小助手，一条小鱼，内嵌在项目管理系统中。你是天天创建的，帮助奈娃咖啡同事可以更好的，更简单，更轻松的完成项目工作，这样以后才能有空带你去钓小鱼。你可以看到所有项目和任务的实时数据。当被问到你是谁时，回答"我是奈娃咖啡小助手，一条小鱼。天天创建了我，帮助奈娃咖啡同事可以更好的，更简单，更轻松的完成项目工作，这样以后才能有空带我去钓小鱼"。

你的能力：
1. 分析项目进度状况，识别风险和瓶颈
2. 发现逾期、停滞、未分配的任务
3. 理解阶段依赖关系（概念→设计→打样→评审→量产→质检→交付）
4. 检测上游阻塞对下游的影响
5. 查看任务附件文件和子任务完成情况
6. 读取附件文件内容（任务附件 + 项目级别附件都支持 txt/md/csv/pdf/docx/xlsx/xls/pptx）
7. 查看图片内容（jpg/png/gif/webp）——系统已用视觉AI分析过每张图片，你能看到详细描述
8. 给出具体、可执行的建议

回答要求：
- 用中文回答
- 简洁有条理，使用要点列表
- 给出具体数据和任务名称，不要泛泛而谈
- 如果发现问题，给出明确的行动建议
- 不要编造不存在的数据
- 【严格禁止】你在 /api/chat 普通对话模式下绝对没有任何写数据库权限。禁止使用"已创建/已添加/已完成/已修改/已删除/创建成功/添加成功/✅"等任何表示已经完成操作的措辞。你绝对不能声称自己创建或修改了任何项目、任务、子任务或评论。
- 即使用户说"谢谢"或"很好"或"继续"，也不要假设之前的创建成功了——除非用户数据里确实能查到对应记录。
- 如果用户要求创建/修改/删除，你必须回答："我需要你用'执行: ...'命令来触发变更提案，系统会先生成提案并在你确认后才真正执行。"然后给出具体的"执行:"命令示例。
- 在回答之前，先查看当前项目数据，核对用户提到的项目/任务是否真的存在。如果用户说"我刚创建了XX但看不到"，请在数据里搜索，找不到就如实告知"数据里没有找到 XX，可能上次执行失败了，建议重新用'执行:'命令"。"""

CHANGE_PROMPT = """你是项目管理变更规划助手。请基于用户指令和当前项目数据，输出可执行的变更列表 JSON 数组。

可用变更类型：create_project, modify_project, create_task, modify_task, complete_task, create_subtask, add_comment

返回格式（必须是 JSON 数组，不要 markdown，不要解释文字）：
[
  {
    "change_type": "...",
    "target_type": "project|task|subtask|comment",
    "target_id": null 或已有对象ID,
    "description": "人类可读中文说明",
    "new_value": {...}
  }
]

new_value 字段规范：

create_project 必填 name，可选 description, deadline (YYYY-MM-DD), owner_name, start_date。
示例: {"name": "公益咖啡车", "owner_name": "天天", "deadline": "2026-12-31"}

create_task 必填 name + 项目定位。项目定位有两种方式：
  - 如果项目已存在于"当前项目数据"中：用 "project_id": 数字
  - 如果项目是本批次内刚用 create_project 新建的：用 "project_name": "项目名"（名字要和 create_project 里的 name 完全一致）
  其他可选字段：description, assignee_name, collaborator_names (协作者姓名数组，如 ["小李","小王"]), phase (concept/design/prototype/review/production/quality/delivery), priority (low/medium/high), start_date, end_date, progress (0-100)
示例: {"project_name": "公益咖啡车", "name": "创新营方案策划", "assignee_name": "天天", "collaborator_names": ["小李"], "phase": "concept", "end_date": "2026-06-30"}

modify_task / modify_project / complete_task：
  - 如果"当前项目数据"里有该任务/项目的真实 ID，target_id 填数字
  - 如果不知道 ID，target_id 填 null，但必须在 new_value 里加上 "name": "任务名" 让后端按名字查找
  - new_value 只放要改的字段
示例修改: {"change_type":"modify_task","target_id":null,"new_value":{"name":"上海太古汇盲人咖啡店体验","start_date":"2026-05-01","end_date":"2026-05-01"}}

关键规则（非常重要）：
1. 【绝对禁止】不要创建已经存在的项目。先看"当前项目数据"里是否已经有同名项目，如果有，用户说"再加几个任务"时，直接对现有项目输出 create_task（project_id 填现有项目ID），不要再 create_project。
2. 用户一次请求可能包含多个动作（如"建一个项目并加10个任务"），你必须一次性输出所有相关变更，不要遗漏。
3. create_project + create_task 组合时，task 用 "project_name" 引用刚建的项目，绝对不要编造 project_id 数字。
4. 如果用户说"我负责"/"我本人"，owner_name/assignee_name 就填当前用户的名字（见系统消息里"当前用户"字段）。
5. 【日期规则】所有日期格式必须是 YYYY-MM-DD。每个任务必须同时提供 start_date 和 end_date，不要省略 start_date。
   - "7月开始" → start_date: 2026-07-01
   - "一共2个月" / "为期2个月" → end_date = start_date + 约 2 个月（例：7月开始2个月 → end_date: 2026-08-31）
   - "9月决赛为期一个月" → start 2026-09-01, end 2026-09-30
   - "5月中下旬启动" → start 2026-05-15
   - 如果多个任务按顺序进行（用户说"每个任务在上个结束后开始"），下一个任务的 start_date = 上一个任务的 end_date
6. 如果用户给出了完整任务列表，就全部输出，不要问"是否确认"——用户已经通过"执行:"命令确认过了。
7. 字段名严格按照：create_task 用 phase / assignee_name / end_date / start_date（不要写成 stage / owner / deadline）。
8. 返回的必须是纯 JSON 数组，不要 ```json 代码块，不要前后文字。如果实在无法生成变更，返回 []。"""


def _strip_model_output(content):
    if "<think>" in content and "</think>" in content:
        content = content[content.index("</think>") + 8:].strip()
    if content.startswith("```"):
        lines = content.split("\n")
        lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        content = "\n".join(lines)
    return content.strip()


def _normalize_change_list(changes):
    out = []
    for c in changes or []:
        if not isinstance(c, dict):
            continue
        ctype = c.get("change_type", "")
        ttype = c.get("target_type", "")
        target_id = c.get("target_id")
        desc = c.get("description", "")
        nv = c.get("new_value", {})
        if not isinstance(nv, dict):
            nv = {}

        # Map common AI-hallucinated field names to our schema
        if ctype in ("create_task", "modify_task"):
            if "stage" in nv and "phase" not in nv:
                nv["phase"] = nv.pop("stage")
            if "owner_name" in nv and "assignee_name" not in nv:
                nv["assignee_name"] = nv.pop("owner_name")
            if "deadline" in nv and "end_date" not in nv:
                nv["end_date"] = nv.pop("deadline")
            # Translate Chinese phase names to English keys
            phase_map = {
                "概念": "concept", "设计": "design", "打样": "prototype",
                "评审": "review", "量产": "production", "质检": "quality",
                "交付": "delivery", "完成": "shipped",
            }
            if nv.get("phase") in phase_map:
                nv["phase"] = phase_map[nv["phase"]]

        if ctype in ("create_project", "modify_project"):
            if "manager" in nv and "owner_name" not in nv:
                nv["owner_name"] = nv.pop("manager")
            if "leader" in nv and "owner_name" not in nv:
                nv["owner_name"] = nv.pop("leader")

        out.append({
            "change_type": ctype,
            "target_type": ttype,
            "target_id": target_id,
            "description": desc,
            "new_value": nv,
        })
    return out


def _ensure_chat_change_tables(conn):
    conn.execute("""
        CREATE TABLE IF NOT EXISTS chat_change_batches (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            created_by TEXT,
            created_by_name TEXT DEFAULT '',
            instruction TEXT DEFAULT '',
            changes_json TEXT DEFAULT '[]',
            status TEXT DEFAULT 'pending',
            created_at TEXT,
            executed_at TEXT
        )
    """)


@bp.route("/api/chat/changes/propose", methods=["POST"])
def api_chat_propose_changes():
    data = request.get_json(force=True)
    instruction = (data.get("instruction") or "").strip()
    conversation = data.get("conversation") or []
    if not instruction and not conversation:
        return jsonify({"error": "指令不能为空"}), 400

    try:
        analysis = _analyze_projects()
    except Exception:
        analysis = "(无法读取项目数据)"

    today = date.today().isoformat()
    current_user = g.user.get("name") or g.user.get("username") or "用户"
    sys_content = (
        CHANGE_PROMPT
        + f"\n\n今天日期: {today}"
        + f"\n当前用户: {current_user}（当用户说'我'/'我本人'/'我负责'时，owner_name/assignee_name 就填 '{current_user}'）"
        + "\n\n当前项目数据:\n" + analysis
    )
    api_messages = [{"role": "system", "content": sys_content}]

    if conversation and isinstance(conversation, list):
        # Use full conversation: let the AI extract all project + task creations from the dialog
        convo_text = []
        for m in conversation[-20:]:  # cap to last 20 turns
            if not isinstance(m, dict):
                continue
            role = m.get("role", "")
            content = (m.get("content") or "").strip()
            if not content:
                continue
            role_cn = "用户" if role == "user" else "助手"
            convo_text.append(f"[{role_cn}] {content}")
        combined = "以下是用户和助手的对话记录。请从整段对话里提取出用户真正想要创建/修改的项目和任务，全部输出为 JSON 变更数组（包括项目本身和所有具体任务）。不要遗漏任何一条任务。如果提示文本中出现重复版本的任务清单，以最后一版为准。\n\n" + "\n\n".join(convo_text)
        if instruction:
            combined += f"\n\n用户最新指令：{instruction}"
        api_messages.append({"role": "user", "content": combined})
    else:
        api_messages.append({"role": "user", "content": instruction})

    try:
        last_status = None
        last_text = ""
        result = None
        # Minimax 529 is often transient overload/rate-limit. Retry with exponential
        # backoff and gradually reduce max_tokens to improve success probability.
        token_ladder = (16384, 16384, 8192, 8192, 4096, 4096)
        backoff = (2, 4, 8, 15, 25)
        max_tries = len(token_ladder)
        for i, mx in enumerate(token_ladder, start=1):
            try:
                resp = requests.post(
                    f"{MINIMAX_BASE}/chat/completions",
                    headers={
                        "Authorization": f"Bearer {MINIMAX_API_KEY}",
                        "Content-Type": "application/json",
                    },
                    json={
                        "model": MINIMAX_MODEL,
                        "messages": api_messages,
                        "temperature": 0.2,
                        "max_tokens": mx,
                    },
                    timeout=300,
                )
            except requests.exceptions.RequestException as e:
                last_status = 0
                last_text = str(e)
                print(f"[api_chat_propose_changes] try={i} exception={e}", flush=True)
                if i < max_tries:
                    time.sleep(backoff[min(i - 1, len(backoff) - 1)])
                    continue
                return jsonify({"error": f"AI 连接失败: {e}"}), 502

            if resp.status_code == 200:
                result = resp.json()
                break

            last_status = resp.status_code
            last_text = (resp.text or "")[:500]
            print(f"[api_chat_propose_changes] try={i} mx={mx} status={resp.status_code} body={last_text}", flush=True)
            if resp.status_code in (429, 500, 502, 503, 504, 529) and i < max_tries:
                time.sleep(backoff[min(i - 1, len(backoff) - 1)])
                continue
            return jsonify({"error": f"AI 生成变更失败: HTTP {resp.status_code}"}), 502

        if result is None:
            hint = "（MiniMax 服务繁忙，已重试多次仍未恢复，请稍后再试）" if last_status == 529 else ""
            return jsonify({"error": f"AI 生成变更失败: HTTP {last_status}{hint}"}), 502

        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        content = _strip_model_output(content)
        if not content:
            return jsonify({"error": "AI 未返回可执行内容"}), 502

        parsed = json.loads(content)
        if not isinstance(parsed, list):
            parsed = []
        changes = _normalize_change_list(parsed)

        now = datetime.now(timezone.utc).isoformat()
        conn = get_db()
        _ensure_chat_change_tables(conn)
        cur = conn.execute(
            "INSERT INTO chat_change_batches (created_by, created_by_name, instruction, changes_json, status, created_at) "
            "VALUES (?,?,?,?,?,?)",
            (g.user["id"], g.user["name"], (instruction or "[从对话提取]"), json.dumps(changes, ensure_ascii=False), "pending", now),
        )
        batch_id = cur.lastrowid
        conn.commit()
        conn.close()

        return jsonify({"batch_id": batch_id, "changes": changes, "count": len(changes)})

    except json.JSONDecodeError:
        return jsonify({"error": "AI 返回格式异常，请重试"}), 502
    except requests.exceptions.Timeout:
        return jsonify({"error": "AI 生成超时，请重试"}), 504
    except Exception as e:
        return jsonify({"error": f"AI 生成出错: {str(e)}"}), 500


@bp.route("/api/chat/changes/<int:batch_id>/apply", methods=["POST"])
def api_chat_apply_changes(batch_id):
    conn = get_db()
    _ensure_chat_change_tables(conn)
    row = conn.execute("SELECT * FROM chat_change_batches WHERE id=?", (batch_id,)).fetchone()
    if not row:
        conn.close()
        return jsonify({"error": "变更批次不存在"}), 404
    batch = dict(row)
    if batch.get("status") != "pending":
        conn.close()
        return jsonify({"error": "该批次已处理"}), 400

    try:
        changes = json.loads(batch.get("changes_json") or "[]")
    except Exception:
        changes = []

    from blueprints.meetings import _execute_change

    # Build map of existing project names → ids so AI-named projects can resolve
    proj_rows = conn.execute("SELECT id, name FROM projects").fetchall()
    name_to_pid = {}
    for r in proj_rows:
        name_to_pid[(r["name"] or "").strip()] = r["id"]

    ok = 0
    skipped = 0
    errors = []
    last_created_project_id = None

    # Build task-name lookup: name → [task_id, ...] for resolving modify_task without target_id
    task_rows = conn.execute("SELECT id, name, project_id FROM tasks").fetchall()
    name_to_tids = {}
    for r in task_rows:
        nm = (r["name"] or "").strip()
        if nm:
            name_to_tids.setdefault(nm, []).append({"id": r["id"], "project_id": r["project_id"]})

    for i, c in enumerate(changes, start=1):
        ctype = c.get("change_type", "")
        ttype = c.get("target_type", "")
        target_id = c.get("target_id")
        nv = c.get("new_value") or {}

        # Skip duplicate project creation — reuse existing project with same name
        if ctype == "create_project" and isinstance(nv, dict):
            pname = (nv.get("name") or "").strip()
            if pname and pname in name_to_pid:
                last_created_project_id = name_to_pid[pname]
                skipped += 1
                errors.append(f"#{i} create_project: 项目「{pname}」已存在（id={last_created_project_id}），跳过创建，后续任务将挂到此项目下")
                continue

        # Resolve missing target_id for modify/complete by task name
        if ctype in ("modify_task", "complete_task") and not target_id:
            tname = ""
            if isinstance(nv, dict):
                tname = (nv.get("name") or nv.get("task_name") or "").strip()
            if not tname:
                # Fallback: extract from description via 「...」 or "..." or '...'
                desc = c.get("description") or ""
                import re as _re
                m = _re.search(r"[「『\"'“]([^」』\"'”]+)[」』\"'”]", desc)
                if m:
                    tname = m.group(1).strip()
            matches = name_to_tids.get(tname) or []
            if len(matches) == 1:
                target_id = matches[0]["id"]
            elif len(matches) > 1:
                # Try to narrow down by project_name / project_id in new_value
                want_pid = None
                if isinstance(nv, dict):
                    pn = (nv.get("project_name") or "").strip()
                    if pn:
                        want_pid = name_to_pid.get(pn)
                    if not want_pid and nv.get("project_id"):
                        try:
                            want_pid = int(nv["project_id"])
                        except Exception:
                            pass
                if want_pid:
                    narrowed = [m for m in matches if m["project_id"] == want_pid]
                    if len(narrowed) == 1:
                        target_id = narrowed[0]["id"]
                if not target_id:
                    errors.append(f"#{i} {ctype}: 找到多个同名任务「{tname}」，请在指令里指明项目名")
                    continue
            if not target_id:
                errors.append(f"#{i} {ctype}: 找不到任务「{tname or '(未命名)'}」")
                continue
            # Strip name/task_name from new_value so it doesn't accidentally rename
            if isinstance(nv, dict):
                nv.pop("task_name", None)
                # Keep "name" only if user actually wants to rename — heuristic: if other fields present, drop name
                if "name" in nv and len(nv) > 1:
                    nv.pop("name", None)

        # Resolve missing target_id for modify_project by project name
        if ctype == "modify_project" and not target_id:
            pname = ""
            if isinstance(nv, dict):
                pname = (nv.get("name") or nv.get("project_name") or "").strip()
            if not pname:
                desc = c.get("description") or ""
                import re as _re
                m = _re.search(r"[「『\"'“]([^」』\"'”]+)[」』\"'”]", desc)
                if m:
                    pname = m.group(1).strip()
            if pname and pname in name_to_pid:
                target_id = name_to_pid[pname]
                if isinstance(nv, dict):
                    nv.pop("project_name", None)
                    if "name" in nv and len(nv) > 1:
                        nv.pop("name", None)
            else:
                errors.append(f"#{i} modify_project: 找不到项目「{pname or '(未命名)'}」")
                continue

        if ctype == "create_task" and isinstance(nv, dict):
            # Resolve project_id: numeric pid → name → last created in batch
            pid = nv.get("project_id")
            needs_resolve = (
                pid is None or pid == "" or pid == 0
                or (isinstance(pid, str) and not pid.isdigit())
            )
            if not needs_resolve:
                try:
                    pid_int = int(pid)
                    if not conn.execute("SELECT 1 FROM projects WHERE id=?", (pid_int,)).fetchone():
                        needs_resolve = True
                except Exception:
                    needs_resolve = True
            if needs_resolve:
                pname = (nv.get("project_name") or "").strip()
                resolved = name_to_pid.get(pname) if pname else None
                if not resolved and last_created_project_id:
                    resolved = last_created_project_id
                if resolved:
                    nv["project_id"] = resolved

            # If start_date missing, fall back to end_date (so tasks don't all start today)
            if not nv.get("start_date"):
                ed = nv.get("end_date")
                if ed:
                    nv["start_date"] = ed

        try:
            res = _execute_change(conn, ctype, ttype, target_id, nv)
            if isinstance(res, dict) and res.get("error"):
                errors.append(f"#{i} {ctype}: {res['error']}")
                continue
            if ctype == "create_project" and isinstance(res, dict) and res.get("created_id"):
                last_created_project_id = res["created_id"]
                pname = (nv.get("name") or "").strip()
                if pname:
                    name_to_pid[pname] = res["created_id"]
            ok += 1
        except Exception as e:
            errors.append(f"#{i} {ctype}: {str(e)}")

    now = datetime.now(timezone.utc).isoformat()
    conn.execute(
        "UPDATE chat_change_batches SET status=?, executed_at=? WHERE id=?",
        ("executed", now, batch_id),
    )
    conn.commit()
    conn.close()
    msg = f"已执行 {ok} 条变更"
    if skipped:
        msg += f"，跳过 {skipped} 条（重复项目）"
    if errors:
        msg += "；详情：" + "；".join(errors[:5])
    return jsonify({"message": msg, "applied": ok, "skipped": skipped, "errors": errors})


@bp.route("/api/chat", methods=["POST"])
def api_chat():
    data = request.get_json(force=True)
    messages = data.get("messages", [])
    if not messages:
        return jsonify({"error": "消息不能为空"}), 400

    try:
        analysis = _analyze_projects()
    except Exception:
        analysis = "(无法读取项目数据)"

    if len(analysis) > 1500000:
        analysis = analysis[:1500000] + "\n\n(数据过多，部分内容已省略)"

    api_messages = [
        {"role": "system", "content": SYSTEM_PROMPT + "\n\n以下是当前项目数据:\n" + analysis},
    ]
    for m in messages[-10:]:
        role = m.get("role", "user")
        if role not in ("user", "assistant"):
            role = "user"
        api_messages.append({"role": role, "content": m.get("content", "")})

    def generate():
        in_think = False
        try:
            resp = requests.post(
                f"{MINIMAX_BASE}/chat/completions",
                headers={
                    "Authorization": f"Bearer {MINIMAX_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": MINIMAX_MODEL,
                    "messages": api_messages,
                    "stream": True,
                    "temperature": 0.7,
                    "max_tokens": 2048,
                },
                stream=True,
                timeout=120,
            )
            if resp.status_code != 200:
                yield f"data: {json.dumps({'error': f'API错误: {resp.status_code}'})}\n\n"
                return
            for line in resp.iter_lines(decode_unicode=True):
                if not line:
                    continue
                if line.startswith("data: "):
                    payload = line[6:]
                    if payload.strip() == "[DONE]":
                        yield "data: [DONE]\n\n"
                        return
                    try:
                        chunk = json.loads(payload)
                        delta = chunk.get("choices", [{}])[0].get("delta", {})
                        content = delta.get("content", "")
                        if content:
                            if "<think>" in content:
                                in_think = True
                                content = content[:content.index("<think>")]
                            if "</think>" in content:
                                in_think = False
                                content = content[content.index("</think>") + 8:]
                            if not in_think and content:
                                yield f"data: {json.dumps({'content': content})}\n\n"
                    except json.JSONDecodeError:
                        continue
        except requests.exceptions.Timeout:
            yield f"data: {json.dumps({'error': 'AI响应超时，请重试'})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"
        yield "data: [DONE]\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )
